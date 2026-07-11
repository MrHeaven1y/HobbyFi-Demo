import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Brand Color Palette (Premium Indigo & Navy Dark Mode / Slate Light Mode)
DARK_BG = RGBColor(15, 23, 42)       # Slate 900 (Dark background for title/conclusion)
LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50 (Light background for content)
PRIMARY = RGBColor(79, 70, 229)      # Indigo 600 (Primary branding color)
SECONDARY = RGBColor(16, 185, 129)   # Emerald 500 (Accent color for features/success)
DARK_TEXT = RGBColor(15, 23, 42)     # Slate 900 for dark headers
LIGHT_TEXT = RGBColor(71, 85, 105)   # Slate 600 for body text
WHITE = RGBColor(255, 255, 255)
CARD_BG = RGBColor(255, 255, 255)    # White background for content cards
CARD_BORDER = RGBColor(226, 232, 240)# Slate 200 for card borders

def set_slide_background(slide, color):
    """Sets the solid background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_header(slide, title_text, is_dark=False):
    """Creates a standard header for content slides."""
    # Top accent band
    band_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(0.4))
    band_shape.fill.solid()
    band_shape.fill.fore_color.rgb = PRIMARY
    band_shape.line.fill.background()
    
    # Title Text Box (Plenty of space to prevent wrap overlap)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(14.4), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE if is_dark else DARK_TEXT

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    """Creates a rounded rectangle container (card) to hold content."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_textbox(slide, left, top, width, height, text="", font_size=18, bold=False, color=LIGHT_TEXT, align=PP_ALIGN.LEFT):
    """Helper to add a standard text box with custom formatting."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = 'Segoe UI'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box

# ==============================================================================
# SLIDE 1: Title Slide (Dark Theme)
# ==============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, DARK_BG)

# Subtle visual shape in background
decor = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10), Inches(0), Inches(6), Inches(9))
decor.fill.solid()
decor.fill.fore_color.rgb = RGBColor(30, 41, 59) # Darker Slate
decor.line.fill.background()
decor.rotation = 180

# Main Title & Subtitles in a single text frame to avoid overlap
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.0), Inches(10), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

p = tf.paragraphs[0]
p.text = "HobbyFi Copilot"
p.font.name = 'Trebuchet MS'
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = WHITE
p.space_after = Pt(8)

p2 = tf.add_paragraph()
p2.text = "AI-Powered CRM & Intelligent Guardrails"
p2.font.name = 'Segoe UI'
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = PRIMARY
p2.space_after = Pt(8)

p3 = tf.add_paragraph()
p3.text = "A Secure, Resilient Copilot for Niche Vendor Portals"
p3.font.name = 'Segoe UI'
p3.font.size = Pt(20)
p3.font.color.rgb = RGBColor(148, 163, 184) # Slate 400

# Presenter Details in a single text frame to avoid overlap
pres_box = slide1.shapes.add_textbox(Inches(1), Inches(5.8), Inches(10), Inches(2.2))
tf_pres = pres_box.text_frame
tf_pres.word_wrap = True
tf_pres.margin_left = tf_pres.margin_top = tf_pres.margin_right = tf_pres.margin_bottom = 0

p_pres = tf_pres.paragraphs[0]
p_pres.text = "Developed & Presented by"
p_pres.font.name = 'Segoe UI'
p_pres.font.size = Pt(14)
p_pres.font.color.rgb = RGBColor(148, 163, 184)
p_pres.space_after = Pt(4)

p_name = tf_pres.add_paragraph()
p_name.text = "Dibayendu Mukherjee" # Fixed spelling from Dibyendu to Dibayendu
p_name.font.name = 'Segoe UI'
p_name.font.size = Pt(22)
p_name.font.bold = True
p_name.font.color.rgb = WHITE
p_name.space_after = Pt(8)

p_desc = tf_pres.add_paragraph()
p_desc.text = "Built rapidly with focus on strict security, resilience under limits, and future-ready architectures."
p_desc.font.name = 'Segoe UI'
p_desc.font.size = Pt(14)
p_desc.font.italic = True
p_desc.font.color.rgb = RGBColor(148, 163, 184)


# ==============================================================================
# SLIDE 2: Product Overview & Core Problem
# ==============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, LIGHT_BG)
create_header(slide2, "The Challenge & The Solution")

# Left Column (The Core Problems - single text frame inside card)
add_card(slide2, Inches(0.8), Inches(1.8), Inches(6.8), Inches(6.4))
tf_prob = slide2.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_prob.word_wrap = True
tf_prob.margin_left = tf_prob.margin_top = tf_prob.margin_right = tf_prob.margin_bottom = 0

p_title1 = tf_prob.paragraphs[0]
p_title1.text = "The Problem Space"
p_title1.font.name = 'Trebuchet MS'
p_title1.font.size = Pt(24)
p_title1.font.bold = True
p_title1.font.color.rgb = DARK_TEXT
p_title1.space_after = Pt(16)

bullets_prob = [
    ("Unsecure Write Mutations", "Traditional AI integration allows the LLM to directly write to the database. A single hallucination or malicious injection could compromise the records."),
    ("Multi-Tenant Leakage Risks", "Portals hold data for multiple vendors. Relying on an LLM to self-restrict queries will fail. Secure scoping must be deterministic."),
    ("Demo & Quota Failures", "Free-tier LLMs (like Gemini) enforce strict quotas. Under active review or presentation, quota exhaustion shouldn't lead to application crashes.")
]
for label, desc in bullets_prob:
    p = tf_prob.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_prob.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)

# Right Column (Our Solution: HobbyFi Copilot - single text frame inside card)
add_card(slide2, Inches(8.4), Inches(1.8), Inches(6.8), Inches(6.4), bg_color=WHITE, border_color=PRIMARY)
tf_sol = slide2.shapes.add_textbox(Inches(8.8), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_sol.word_wrap = True
tf_sol.margin_left = tf_sol.margin_top = tf_sol.margin_right = tf_sol.margin_bottom = 0

p_title2 = tf_sol.paragraphs[0]
p_title2.text = "The Copilot Solution"
p_title2.font.name = 'Trebuchet MS'
p_title2.font.size = Pt(24)
p_title2.font.bold = True
p_title2.font.color.rgb = PRIMARY
p_title2.space_after = Pt(16)

bullets_sol = [
    ("Decoupled Execution Architecture", "Treats LLM purely as a natural-language reasoning agent that outputs structured payloads, rather than allowing direct database writes."),
    ("Deterministic Scoping Rules", "Backend layers intercept and overwrite all LLM requests with the cryptographically authenticated vendor scope."),
    ("Resilient Fallback Design", "A 3-tier fail-safe chain that degrades gracefully to local offline models or deterministic responses instead of raising a 500 crash.")
]
for label, desc in bullets_sol:
    p = tf_sol.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_sol.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)


# ==============================================================================
# SLIDE 3: Accurate Tech Stack
# ==============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, LIGHT_BG)
create_header(slide3, "The Production Technology Stack")

# Grid of 4 cards
tech_items = [
    ("Backend & Orchestration", ["FastAPI (High-performance API)", "LangGraph (Cyclic state agent)", "LangChain (Tool-calling SDK)"], Inches(0.8), Inches(1.8)),
    ("Database & Vector Storage", ["PostgreSQL 16 (Source of truth)", "pgvector extension (Vector DB)", "SQLAlchemy (Strict ORM layer)"], Inches(8.2), Inches(1.8)),
    ("AI Models & Inference", ["Gemini 2.5 Flash (Primary LLM)", "Ollama & Llama 3.2 (Local fallback)", "Hugging Face Inference (Embeddings)"], Inches(0.8), Inches(4.9)),
    ("Lightweight Modern UI", ["Jinja2 templates (Server rendered)", "Alpine.js (Reactive front-end state)", "HTMX (Clean async HTML requests)"], Inches(8.2), Inches(4.9))
]

for title, bullets, left, top in tech_items:
    add_card(slide3, left, top, Inches(7.0), Inches(2.7))
    tf_tech = slide3.shapes.add_textbox(left + Inches(0.4), top + Inches(0.3), Inches(6.2), Inches(2.1)).text_frame
    tf_tech.word_wrap = True
    tf_tech.margin_left = tf_tech.margin_top = tf_tech.margin_right = tf_tech.margin_bottom = 0
    
    p_title = tf_tech.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY
    p_title.space_after = Pt(12)
    
    for item in bullets:
        p = tf_tech.add_paragraph()
        p.text = f"• {item}"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(15)
        p.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(6)


# ==============================================================================
# SLIDE 4: Architecture Flow & Decentralized Reasoning
# ==============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, LIGHT_BG)
create_header(slide4, "Core Request Lifecycle & Architecture")

# Visual Box representing the chain
flow_steps = [
    ("1. Browser UI", "HTMX sends chat input", Inches(0.8), Inches(2.2), Inches(2.6)),
    ("2. FastAPI Layer", "Auth middleware scopes Vendor ID", Inches(3.8), Inches(2.2), Inches(2.6)),
    ("3. Copilot Service", "Injects DB & active memory window", Inches(6.8), Inches(2.2), Inches(2.6)),
    ("4. LangGraph Agent", "Drives prompt & Gemini logic loop", Inches(9.8), Inches(2.2), Inches(2.6)),
    ("5. CRM Tools / DB", "Executes read / drafts pending writes", Inches(12.8), Inches(2.2), Inches(2.4))
]

for title, desc, left, top, width in flow_steps:
    add_card(slide4, left, top, width, Inches(2.0), border_color=PRIMARY)
    
    tf_flow = slide4.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), Inches(1.6)).text_frame
    tf_flow.word_wrap = True
    tf_flow.margin_left = tf_flow.margin_top = tf_flow.margin_right = tf_flow.margin_bottom = 0
    
    p_flow_title = tf_flow.paragraphs[0]
    p_flow_title.text = title
    p_flow_title.alignment = PP_ALIGN.CENTER
    p_flow_title.font.name = 'Segoe UI'
    p_flow_title.font.size = Pt(18)
    p_flow_title.font.bold = True
    p_flow_title.font.color.rgb = PRIMARY
    p_flow_title.space_after = Pt(6)
    
    p_flow_desc = tf_flow.add_paragraph()
    p_flow_desc.text = desc
    p_flow_desc.alignment = PP_ALIGN.CENTER
    p_flow_desc.font.name = 'Segoe UI'
    p_flow_desc.font.size = Pt(13)
    p_flow_desc.font.color.rgb = LIGHT_TEXT

# Core Principle Box below
add_card(slide4, Inches(0.8), Inches(4.7), Inches(14.4), Inches(3.2), bg_color=DARK_BG, border_color=None)
tf_phil = slide4.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(13.6), Inches(2.6)).text_frame
tf_phil.word_wrap = True
tf_phil.margin_left = tf_phil.margin_top = tf_phil.margin_right = tf_phil.margin_bottom = 0

p_phil_title = tf_phil.paragraphs[0]
p_phil_title.text = "Core Architecture Philosophy"
p_phil_title.font.name = 'Trebuchet MS'
p_phil_title.font.size = Pt(20)
p_phil_title.font.bold = True
p_phil_title.font.color.rgb = SECONDARY
p_phil_title.space_after = Pt(10)

principles = [
    ("The Model Reasons", "The LLM generates tool arguments and evaluates user intent. It acts solely as a natural-language router."),
    ("The Backend Enforces", "The FastAPI and SQLAlchemy layers govern multi-tenant scoping and transaction security. Security is hardcoded, never prompted."),
    ("The Database Secures", "Audit logs and pgvector records preserve full data isolation and track pending/executed transactions."),
    ("The Vendor Confirms", "Human-in-the-Loop approval guarantees that no business state is altered without explicit confirmation.")
]
for title, desc in principles:
    p = tf_phil.add_paragraph()
    p.text = f"{title}: "
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE
    
    run = p.add_run()
    run.text = desc
    run.font.bold = False
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(203, 213, 225) # Slate 300
    p.space_after = Pt(6)


# ==============================================================================
# SLIDE 5: Security Guardrails & Tenant Isolation
# ==============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5, LIGHT_BG)
create_header(slide5, "Security Guardrails & Multi-Tenant Isolation")

# Left: Cryptographic Tenant Isolation (single textframe)
add_card(slide5, Inches(0.8), Inches(1.8), Inches(6.8), Inches(6.4))
tf_guard1 = slide5.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_guard1.word_wrap = True
tf_guard1.margin_left = tf_guard1.margin_top = tf_guard1.margin_right = tf_guard1.margin_bottom = 0

p_g1_title = tf_guard1.paragraphs[0]
p_g1_title.text = "Tenant Isolation (Zero Trust)"
p_g1_title.font.name = 'Trebuchet MS'
p_g1_title.font.size = Pt(24)
p_g1_title.font.bold = True
p_g1_title.font.color.rgb = DARK_TEXT
p_g1_title.space_after = Pt(16)

bullets_guard1 = [
    ("Middleware Authentication", "The `VendorAuthMiddleware` extracts and cryptographically validates the `X-Vendor-ID` header from incoming requests at the API gateway."),
    ("Tool-Level Override", "Even if a hallucinating LLM requests data for another vendor (e.g. `v_67890_xyz`), the Python CRM tool ignores the argument and overrides it with the verified vendor ID from context."),
    ("No Cross-Vendor Leaks", "Strict SQL filters are applied at the query generation stage. Tested successfully to reject cross-vendor read/write requests.")
]
for label, desc in bullets_guard1:
    p = tf_guard1.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_guard1.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)

# Right: Human-In-The-Loop Approvals (single textframe)
add_card(slide5, Inches(8.4), Inches(1.8), Inches(6.8), Inches(6.4))
tf_guard2 = slide5.shapes.add_textbox(Inches(8.8), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_guard2.word_wrap = True
tf_guard2.margin_left = tf_guard2.margin_top = tf_guard2.margin_right = tf_guard2.margin_bottom = 0

p_g2_title = tf_guard2.paragraphs[0]
p_g2_title.text = "Human-In-The-Loop (HITL)"
p_g2_title.font.name = 'Trebuchet MS'
p_g2_title.font.size = Pt(24)
p_g2_title.font.bold = True
p_g2_title.font.color.rgb = PRIMARY
p_g2_title.space_after = Pt(16)

bullets_guard2 = [
    ("Non-Mutating Write Tools", "Write tools (`update_membership`, `extend_trial`) are prohibited from writing directly. They only draft a pending `audit_logs` entry."),
    ("Pending Audit Trail", "Drafted actions are stored as JSON payloads (status: `pending`) mapped to the authenticated vendor, returning a secure ID to the UI."),
    ("Explicit Authorization", "The user reviews a structured card in the UI. When they click 'Approve', a separate API endpoint validates vendor ownership and executes the transaction.")
]
for label, desc in bullets_guard2:
    p = tf_guard2.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_guard2.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)


# ==============================================================================
# SLIDE 6: Resilient Fallback Mechanics
# ==============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6, LIGHT_BG)
create_header(slide6, "Resilient 3-Tier Fallback Mechanism")

# Left Column (The Fail-Safe Chain - single textframe)
add_card(slide6, Inches(0.8), Inches(1.8), Inches(6.8), Inches(6.4))
tf_fail = slide6.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_fail.word_wrap = True
tf_fail.margin_left = tf_fail.margin_top = tf_fail.margin_right = tf_fail.margin_bottom = 0

p_fail_title = tf_fail.paragraphs[0]
p_fail_title.text = "The Resiliency Chain"
p_fail_title.font.name = 'Trebuchet MS'
p_fail_title.font.size = Pt(24)
p_fail_title.font.bold = True
p_fail_title.font.color.rgb = DARK_TEXT
p_fail_title.space_after = Pt(16)

bullets_fail = [
    ("Tier 1: Gemini 2.5 Flash (Primary)", "Uses full LangGraph state machine orchestration and tool-calling models for standard, complex natural language routing."),
    ("Tier 2: Local Ollama Model (Secondary)", "If Gemini quota exhausts, the backend automatically redirects to a configured local model (e.g. `llama3.2`) with injected context."),
    ("Tier 3: Deterministic Fallback (Tertiary)", "If local model is also offline, it relies on structured, deterministic CRM responses so the system never breaks during demo review.")
]
for label, desc in bullets_fail:
    p = tf_fail.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_fail.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)

# Right Column (Why This Matters & Implementation - single textframe)
add_card(slide6, Inches(8.4), Inches(1.8), Inches(6.8), Inches(6.4), bg_color=WHITE, border_color=SECONDARY)
tf_details = slide6.shapes.add_textbox(Inches(8.8), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_details.word_wrap = True
tf_details.margin_left = tf_details.margin_top = tf_details.margin_right = tf_details.margin_bottom = 0

p_det_title = tf_details.paragraphs[0]
p_det_title.text = "Implementation Details"
p_det_title.font.name = 'Trebuchet MS'
p_det_title.font.size = Pt(24)
p_det_title.font.bold = True
p_det_title.font.color.rgb = SECONDARY
p_det_title.space_after = Pt(16)

bullets_details = [
    ("Transparent UI Badging", "The interface displays the active mode ('Gemini', 'Local Model', or 'Fallback') and remaining LLM budget for full developer visibility."),
    ("Staged Loading Alerts", "Warns users when the organization quota is low and fallbacks are executing to manage expectations during reviews."),
    ("Runtime Event Auditing", "Every transition to fallback modes is logged in the `runtime_events` table for operational visibility and debugging.")
]
for label, desc in bullets_details:
    p = tf_details.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_details.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)


# ==============================================================================
# SLIDE 7: The Copilot Interface in Action (Screenshots)
# ==============================================================================
slide_img = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide_img, LIGHT_BG)
create_header(slide_img, "The Copilot Interface in Action")

# Left Image: Approved & Executed Write Action
pic_path_left = "Pictures/chrome_0TWUGEPASU.png"
if os.path.exists(pic_path_left):
    # Image frame
    add_card(slide_img, Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.2), bg_color=WHITE, border_color=PRIMARY)
    slide_img.shapes.add_picture(pic_path_left, Inches(1.0), Inches(2.0), width=Inches(6.4), height=Inches(3.8))
    
    # Caption box
    tf_cap_left = slide_img.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(6.8), Inches(1.8)).text_frame
    tf_cap_left.word_wrap = True
    p_cap_left = tf_cap_left.paragraphs[0]
    p_cap_left.text = "Human-in-the-Loop Approvals"
    p_cap_left.font.name = 'Segoe UI'
    p_cap_left.font.size = Pt(18)
    p_cap_left.font.bold = True
    p_cap_left.font.color.rgb = PRIMARY
    p_cap_left.space_after = Pt(4)
    
    p_cap_left_desc = tf_cap_left.add_paragraph()
    p_cap_left_desc.text = "Demonstrates a membership write mutation that requires approval, is approved via the UI card, and then successfully executed."
    p_cap_left_desc.font.name = 'Segoe UI'
    p_cap_left_desc.font.size = Pt(14)
    p_cap_left_desc.font.color.rgb = LIGHT_TEXT

# Right Image: Guardrail Testing & Cross-Vendor Refusal
pic_path_right = "Pictures/chrome_N0ZVgJtSap.png"
if os.path.exists(pic_path_right):
    # Image frame
    add_card(slide_img, Inches(8.4), Inches(1.8), Inches(6.8), Inches(4.2), bg_color=WHITE, border_color=PRIMARY)
    slide_img.shapes.add_picture(pic_path_right, Inches(8.6), Inches(2.0), width=Inches(6.4), height=Inches(3.8))
    
    # Caption box
    tf_cap_right = slide_img.shapes.add_textbox(Inches(8.4), Inches(6.2), Inches(6.8), Inches(1.8)).text_frame
    tf_cap_right.word_wrap = True
    p_cap_right = tf_cap_right.paragraphs[0]
    p_cap_right.text = "Multi-Tenant Guardrails"
    p_cap_right.font.name = 'Segoe UI'
    p_cap_right.font.size = Pt(18)
    p_cap_right.font.bold = True
    p_cap_right.font.color.rgb = PRIMARY
    p_cap_right.space_after = Pt(4)
    
    p_cap_right_desc = tf_cap_right.add_paragraph()
    p_cap_right_desc.text = "Testing cross-vendor data queries (e.g. database dump request). The copilot securely refuses access, preventing information leaks."
    p_cap_right_desc.font.name = 'Segoe UI'
    p_cap_right_desc.font.size = Pt(14)
    p_cap_right_desc.font.color.rgb = LIGHT_TEXT


# ==============================================================================
# SLIDE 8: Ingestion Pipeline Feature
# ==============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7, LIGHT_BG)
create_header(slide7, "New Feature: Background Document Ingestion Pipeline")

# Big card at the top (single textframe)
add_card(slide7, Inches(0.8), Inches(1.8), Inches(14.4), Inches(2.2))
tf_ing = slide7.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(13.6), Inches(1.6)).text_frame
tf_ing.word_wrap = True
tf_ing.margin_left = tf_ing.margin_top = tf_ing.margin_right = tf_ing.margin_bottom = 0

p_ing_title = tf_ing.paragraphs[0]
p_ing_title.text = "Asynchronous Support Document Ingestion"
p_ing_title.font.name = 'Segoe UI'
p_ing_title.font.size = Pt(22)
p_ing_title.font.bold = True
p_ing_title.font.color.rgb = PRIMARY
p_ing_title.space_after = Pt(8)

p_ing_desc = tf_ing.add_paragraph()
p_ing_desc.text = "In this update, I built and integrated a background document ingestion pipeline (`POST /api/v1/documents/upload`). While the current phase of the Copilot focuses strictly on structured CRM data, this new feature establishes the foundational architecture for future Unstructured RAG capabilities."
p_ing_desc.font.name = 'Segoe UI'
p_ing_desc.font.size = Pt(16)
p_ing_desc.font.color.rgb = LIGHT_TEXT

# Flow steps below
add_textbox(slide7, Inches(0.8), Inches(4.3), Inches(14.4), Inches(0.4), "How the Ingestion Pipeline Operates", 20, True, DARK_TEXT)

pipeline_steps = [
    ("1. FastAPI Endpoint", "Accepts multipart file upload and hands off to FastAPI background task.", Inches(0.8), Inches(4.8), Inches(3.4)),
    ("2. Extraction & Clean", "Saves temporarily, extracts text (OCR/Parser), and strips extra whitespace.", Inches(4.4), Inches(4.8), Inches(3.4)),
    ("3. Chunking & Overlap", "Chunks clean text using sliding thresholds (15 lines per chunk, 5 lines overlap).", Inches(8.0), Inches(4.8), Inches(3.4)),
    ("4. pgvector Embedding", "Embeds chunk text and persists to PostgreSQL with pgvector embeddings.", Inches(11.6), Inches(4.8), Inches(3.6))
]

for title, desc, left, top, width in pipeline_steps:
    add_card(slide7, left, top, width, Inches(3.0), border_color=PRIMARY)
    
    tf_pipe = slide7.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), width - Inches(0.4), Inches(2.4)).text_frame
    tf_pipe.word_wrap = True
    tf_pipe.margin_left = tf_pipe.margin_top = tf_pipe.margin_right = tf_pipe.margin_bottom = 0
    
    p_pipe_title = tf_pipe.paragraphs[0]
    p_pipe_title.text = title
    p_pipe_title.font.name = 'Segoe UI'
    p_pipe_title.font.size = Pt(16)
    p_pipe_title.font.bold = True
    p_pipe_title.font.color.rgb = PRIMARY
    p_pipe_title.space_after = Pt(8)
    
    p_pipe_desc = tf_pipe.add_paragraph()
    p_pipe_desc.text = desc
    p_pipe_desc.font.name = 'Segoe UI'
    p_pipe_desc.font.size = Pt(13)
    p_pipe_desc.font.color.rgb = LIGHT_TEXT


# ==============================================================================
# SLIDE 8: Memory & Resource Optimizations (Render Constraints)
# ==============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide8, LIGHT_BG)
create_header(slide8, "DevOps Constraints & Performance Optimizations")

# Left Column (Hugging Face API Optimization - single textframe)
add_card(slide8, Inches(0.8), Inches(1.8), Inches(6.8), Inches(6.4), bg_color=WHITE, border_color=PRIMARY)
tf_oom = slide8.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_oom.word_wrap = True
tf_oom.margin_left = tf_oom.margin_top = tf_oom.margin_right = tf_oom.margin_bottom = 0

p_oom_title = tf_oom.paragraphs[0]
p_oom_title.text = "Solving the Render OOM Problem"
p_oom_title.font.name = 'Trebuchet MS'
p_oom_title.font.size = Pt(24)
p_oom_title.font.bold = True
p_oom_title.font.color.rgb = PRIMARY
p_oom_title.space_after = Pt(16)

bullets_oom = [
    ("The Constraint: 512MB RAM Limit", "Free hosting services like Render enforce a strict 512MB memory ceiling. Loading PyTorch and local models inside the app led to instant Out-Of-Memory (OOM) crashes."),
    ("The Design Shift: Serverless APIs", "I shifted from local `sentence-transformers` library to the Hugging Face Serverless Inference API to compute embeddings."),
    ("The Result: Near-Zero Memory", "Embedding calculations are offloaded entirely. Memory footprints dropped to near-zero, while keeping the 384-dimensional vector database mapping 100% compatible.")
]
for label, desc in bullets_oom:
    p = tf_oom.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_oom.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)

# Right Column (Sliding Memory window - single textframe)
add_card(slide8, Inches(8.4), Inches(1.8), Inches(6.8), Inches(6.4))
tf_mem = slide8.shapes.add_textbox(Inches(8.8), Inches(2.1), Inches(6.0), Inches(5.8)).text_frame
tf_mem.word_wrap = True
tf_mem.margin_left = tf_mem.margin_top = tf_mem.margin_right = tf_mem.margin_bottom = 0

p_mem_title = tf_mem.paragraphs[0]
p_mem_title.text = "Intelligent Session Memory"
p_mem_title.font.name = 'Trebuchet MS'
p_mem_title.font.size = Pt(24)
p_mem_title.font.bold = True
p_mem_title.font.color.rgb = DARK_TEXT
p_mem_title.space_after = Pt(16)

bullets_mem = [
    ("PostgreSQL-Backed Memory", "Conversation states are serialized and loaded directly from PostgreSQL databases based on persistent conversation UUIDs."),
    ("Sliding Window Pruning", "Coded with an active sliding window of the last 10 messages. Older history is dynamically pruned from active contexts during requests."),
    ("Optimizing Latency & Accuracy", "Capping history prevents model hallucinations, enforces clean conversation flows, and limits token usage to maintain fast API response times.")
]
for label, desc in bullets_mem:
    p = tf_mem.add_paragraph()
    p.text = f"•  {label}:"
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(2)
    
    p_desc = tf_mem.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = LIGHT_TEXT
    p_desc.space_after = Pt(14)


# ==============================================================================
# SLIDE 9: Production Roadmap
# ==============================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide9, LIGHT_BG)
create_header(slide9, "Future Production-Ready Roadmap")

# 3 horizontal timeline items
timeline_items = [
    ("Phase 1: Authorization Hardening", [
        "Replace the current simulated X-Vendor-ID header auth with secure JWT or OAuth 2.0 flows.",
        "Implement Role-Based Access Control (RBAC): Admin roles auto-execute, while Staff roles route to Human-in-the-Loop approvals."
    ], Inches(0.8)),
    ("Phase 2: UX & Stream Upgrades", [
        "Upgrade current REST endpoints to WebSockets to support real-time LLM token streaming.",
        "Add rich UI approval card breakdowns to preview exact user/game database details prior to execution."
    ], Inches(5.8)),
    ("Phase 3: Database & CI Hardening", [
        "Transition startup-based database creation to version-controlled Alembic migrations.",
        "Build a CI pipeline that runs safety tests automatically on pull requests."
    ], Inches(10.8))
]

for title, points, left in timeline_items:
    add_card(slide9, left, Inches(1.8), Inches(4.4), Inches(6.4))
    
    tf_timeline = slide9.shapes.add_textbox(left + Inches(0.3), Inches(2.1), Inches(3.8), Inches(5.8)).text_frame
    tf_timeline.word_wrap = True
    tf_timeline.margin_left = tf_timeline.margin_top = tf_timeline.margin_right = tf_timeline.margin_bottom = 0
    
    p_time_title = tf_timeline.paragraphs[0]
    p_time_title.text = title
    p_time_title.font.name = 'Segoe UI'
    p_time_title.font.size = Pt(20)
    p_time_title.font.bold = True
    p_time_title.font.color.rgb = PRIMARY
    p_time_title.space_after = Pt(14)
    
    for pt in points:
        p = tf_timeline.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(16)


# ==============================================================================
# SLIDE 10: Conclusion & Contacts (Dark Theme)
# ==============================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide10, DARK_BG)

# Title Thank You (single textframe)
title_box = slide10.shapes.add_textbox(Inches(1), Inches(1.2), Inches(14), Inches(1.6))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

p = tf.paragraphs[0]
p.text = "Thank You & Product Demo"
p.alignment = PP_ALIGN.CENTER
p.font.name = 'Trebuchet MS'
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.space_after = Pt(6)

p_sub = tf.add_paragraph()
p_sub.text = "Open to questions & code reviews"
p_sub.alignment = PP_ALIGN.CENTER
p_sub.font.name = 'Segoe UI'
p_sub.font.size = Pt(20)
p_sub.font.color.rgb = SECONDARY

# Main container for credentials
add_card(slide10, Inches(2.5), Inches(3.2), Inches(11.0), Inches(4.5), bg_color=RGBColor(30, 41, 59), border_color=None)

tf_contact = slide10.shapes.add_textbox(Inches(3.0), Inches(3.6), Inches(10.0), Inches(3.8)).text_frame
tf_contact.word_wrap = True
tf_contact.margin_left = tf_contact.margin_top = tf_contact.margin_right = tf_contact.margin_bottom = 0

contact_info = [
    ("Live Application Demo", "https://hobbyfi-copilot-wlbk.onrender.com"),
    ("GitHub Repository", "https://github.com/MrHeaven1y"),
    ("LinkedIn Profile", "www.linkedin.com/in/dibayendu-mukherjee-bb897b267"),
    ("Email Address", "dibyendumukherjee916@gmail.com")
]

for label, value in contact_info:
    p = tf_contact.add_paragraph() if tf_contact.text else tf_contact.paragraphs[0]
    p.text = f"{label}: "
    p.font.name = 'Segoe UI'
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = SECONDARY
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = value
    run.font.bold = False
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    p.space_after = Pt(20)

# Save presentation to multiple output filenames in the repo to replace all outdated versions
paths = ["HobbyFi_Presentation.pptx", "HobbyFi_Product_Pitch.pptx", "HobbyFi_Presentation_v2.pptx"]
for path in paths:
    prs.save(path)
    print(f"Presentation saved successfully to: {path}")
